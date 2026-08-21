"""
Slide deck summarising the CURRENT pipeline: which method does what, which
script produces it, and where its outputs land.

Companion to documentation_and_summaries/README.md -- same content, slide form.
Covers only the live pipeline; the retired `A` / trend method gets one closing
line, not an explanation.

Built on documentation_and_summaries/Template.pptx, whose canvas is 10 x 5.62 in
-- much smaller than the PowerPoint default. Font sizes are set for THAT canvas;
do not copy them to a 13.3 in deck. Uses the template's own layouts (#3 branded
title, #14 full-width single column) so the branding and type come out native.

Run in the `sup3r` conda env (the one with python-pptx):
    conda activate sup3r
    python make_pipeline_overview_deck.py

Outputs:
  * documentation_and_summaries/wildfire_pipeline_overview.pptx
"""

import os
from pptx import Presentation
from pptx.util import Inches, Pt

DOCS = ("/projects/alcaps/bfuchs/wildfire_risk_changes_capacity"
        "/documentation_and_summaries")
TEMPLATE = os.path.join(DOCS, "Template.pptx")
OUT = os.path.join(DOCS, "wildfire_pipeline_overview.pptx")

LAYOUT_TITLE = 3      # Title Slide - Branded 02 - Dune
LAYOUT_BODY = 14      # Simple Slide 3 - Text, 1 column  (body 9.0 x 4.0 @0.5,1.2)

TITLE = "Wildfire risk → transmission cost"
SUBTITLE = "What each step does  ·  TVA + Southern California  ·  August 2026"

# (title, [(text, level)]) -- level 0 = main point, 1 = sub-detail
SLIDES = [
    ("\U0001F525  The pipeline", [
        ("Three methods, run in order — each feeds the next.", 0),
        ("①  HISTORICAL RISK ZONES  →  sets the class cutoffs", 0),
        ("②  PROJECTED RISK ZONES  →  same cutoffs, future fire weather", 0),
        ("③  COST PENALTY  →  prices that risk onto transmission routes", 0),
        ("Read ① first: it produces the cutoffs ② and ③ are both "
         "measured against.", 0),
        ("⚠  Fuel is not in any of these yet — see the last slide.", 0),
    ]),

    ("①②  Risk zones — historical, then projected", [
        ("METHOD  Classify each grid cell by its p98 FWI — the fire-weather "
         "severity its worst 2% of days reach. Absolute severity, not a trend.", 0),
        ("non-land cells dropped first; each region gets its own cutoffs from "
         "the 50th / 75th / 90th percentile cell within it", 1),
        ("CUTOFFS   TVA  35.1 / 38.8 / 40.8      SoCal  98.6 / 111.7 / 125.2", 0),
        ("PROJECTED  Recompute the same p98 over 2025–2059, classify with "
         "those historical cutoffs HELD FIXED.", 0),
        ("re-deriving them would force the same 50/25/15/10 split and erase "
         "the projection signal — letting the shares move is the result", 1),
        ("RESULT   TVA high-risk land   10.0%  →  24.9%", 0),
        ("regional_absolute_risk_maps.py  ·  future_fwi_projected_hazard.py"
         "   →  outputs/risk_historical/ , risk_future/", 1),
    ]),

    ("③  ⚡  Cost penalty on transmission routes", [
        ("METHOD  Multiply every 90 m transmission cost cell by the multiplier "
         "of the risk class it sits in, then re-cost the 1,509 existing TVA "
         "routes on the penalised surface.", 0),
        ("none ×1.0      low ×1.1      medium ×1.3      high ×1.5", 0),
        ("“none” is ×1.0, not ×0 — a zero multiplier would make "
         "no-risk ground free and collapse every route onto it", 1),
        ("RESULT   system total  +10.74% on historical zones  →  +20.17% on "
         "projected zones;   median route  +17.68%", 0),
        ("transmission_cost_risk_penalty.py  ·  route_risk_cost_analysis.py"
         "   →  outputs/cost_penalty_historical/ , cost_penalty_future/", 1),
        ("Each script has one switch (FIELD_SOURCE / RISK_SOURCE) routing both "
         "its inputs and outputs, so the two vintages never overwrite.", 0),
    ]),

    ("⚠  What is not in here yet", [
        ("FUEL.  FWI is a fire WEATHER index with no fuel term — it cannot "
         "see that the eastern TVA plateau carries far more fuel than the "
         "agricultural west, so the zones still concentrate in the west.", 0),
        ("fuel and fire weather are near-independent across TVA (rank corr "
         "−0.04) and run in OPPOSITE directions with longitude: "
         "+0.21 fuel vs −0.52 fire weather", 1),
        ("LANDFIRE FBFM40 will redefine the zones with fuel included; "
         "landfire_fuel_composition.py is written, not yet run", 1),
        ("ALL NUMBERS ABOVE ARE PROVISIONAL and will move once fuel is in.", 0),
        ("An earlier trend-based method (the “A” metric) was tried and "
         "retired — documented in documentation_and_summaries/archive/.", 0),
    ]),
]


def strip_slides(prs):
    """Template ships with example slides; remove them."""
    ids = prs.slides._sldIdLst
    for sld in list(ids):
        rId = sld.get("{http://schemas.openxmlformats.org/officeDocument/"
                      "2006/relationships}id")
        prs.part.drop_rel(rId)
        ids.remove(sld)


def main():
    prs = Presentation(TEMPLATE)
    strip_slides(prs)

    s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_TITLE])
    s.shapes.title.text = TITLE
    for ph in s.placeholders:
        if ph.placeholder_format.idx == 11:
            ph.text = SUBTITLE

    for title, lines in SLIDES:
        s = prs.slides.add_slide(prs.slide_layouts[LAYOUT_BODY])
        s.shapes.title.text = title
        body = next(ph for ph in s.placeholders
                    if ph.placeholder_format.idx == 10)
        tf = body.text_frame
        tf.word_wrap = True
        for i, (text, lvl) in enumerate(lines):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.text = text
            p.level = lvl
            # sized for the 10 in canvas; spacing keeps the block from
            # bunching at the top of the 4.0 in body
            p.font.size = Pt(12) if lvl == 0 else Pt(10.5)
            p.space_after = Pt(9) if lvl == 0 else Pt(6)

    prs.save(OUT)
    print(f"  saved -> {OUT}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
