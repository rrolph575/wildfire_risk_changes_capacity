"""
Build the summary slide deck for the FUTURE (projected) method.

Walks the method from data inputs through to the cost penalty, and embeds the
explainer figures from outputs/cost_penalty_future/.

Every number in here is taken from a run of the pipeline, not estimated. If the
pipeline is rerun with different thresholds or multipliers, update the numbers
in SLIDES below.

Run in the `sup3r` conda env (the one with python-pptx):
    conda activate sup3r
    python make_future_method_deck.py

Outputs:
  * outputs/cost_penalty_future/wildfire_future_method_summary.pptx
"""

import os

from pptx import Presentation
from pptx.util import Inches, Pt

PROJ = "/projects/alcaps/bfuchs/wildfire_risk_changes_capacity"
RISK_FUT = os.path.join(PROJ, "outputs", "risk_future")
COST_FUT = os.path.join(PROJ, "outputs", "cost_penalty_future")
OUT = os.path.join(COST_FUT, "wildfire_future_method_summary.pptx")

TITLE = "Projected wildfire risk and transmission cost penalty"
SUBTITLE = ("TVA coverage area  ·  Sup3rCC ssp245 2025-2059  ·  "
            "methods summary\nAugust 2026")

# (title, [bullets], image path or None)
SLIDES = [
    ("1. Data inputs", [
        "Sup3rCC daily FWI (`fwi_tc`) — Canadian Fire Weather Index, "
        "dimensionless, trend-corrected for humidity",
        "    · ~4 km CONUS grid, 2,300,000 cells",
        "    · historical run 2000–2014 (15 yr)  ·  ssp245 run 2025–2059 "
        "(35 yr, ~118 GB read)",
        "USFS Wildfire Hazard Potential, 90 m — used here only as a "
        "burnable / non-burnable land mask",
        "TVA least-cost-path cost surface — 90 m, EPSG:5070, cost to build "
        "through each cell",
        "TVA routes — 1,509 routed lines between 57 buses, with their "
        "as-built costs",
    ], None),

    ("2. Risk classes from absolute fire-weather severity", [
        "Each cell's p98 FWI = the level reached on its worst 2% of days "
        "(~7 days/yr). An absolute severity measure, not a trend.",
        "Cells not on land (ocean, outside CONUS) are dropped first — they "
        "cannot burn, and including them lets water set the thresholds.",
        "Each region gets its own cutoffs, from the 50th / 75th / 90th "
        "percentile cell within that region:",
        "        TVA      low 35.1     medium 38.8     high 40.8",
        "        SoCal    low 98.6     medium 111.7    high 125.2",
        "Regions differ so much that one shared cutoff would leave TVA with "
        "no high cells at all: TVA's worst cell (47.4) sits below SoCal's "
        "25th-percentile cell.",
    ], None),

    ("3. Projecting the field forward", [
        "Recomputed the same quantity — p98 FWI per cell — over the ssp245 "
        "future window 2025–2059 (batch job, 14 min 46 s).",
        "THE KEY CHOICE: the cutoffs are held FIXED at the historical values, "
        "not re-derived from the future field.",
        "    · re-deriving would force the same 50/25/15/10 split by "
        "construction, and warming is roughly uniform in space — the map "
        "would barely change and the projection signal would vanish",
        "    · holding them fixed lets the class shares move, and that "
        "movement IS the result",
        "Result — share of TVA land cells:",
        "        none  50.0% → 29.2%        medium  15.0% → 19.2%",
        "        low   25.0% → 26.7%        high    10.0% → 24.9%",
        "TVA's high-risk area roughly 2.5×; SoCal's roughly 1.7×.",
    ], None),

    ("Projected risk classes, TVA (2025–2059)", [],
     os.path.join(RISK_FUT, "abs_risk_future_p98_tva.png")),

    ("4. Turning risk into a cost penalty", [
        "Every 90 m cost cell is multiplied by the class of the ~2.8 km risk "
        "cell containing it:",
        "        none ×1.0      low ×1.1      medium ×1.3      high ×1.5",
        "\"none\" is ×1.0, not ×0 — a zero multiplier would make no-risk "
        "ground free rather than unpenalised, and the router would collapse "
        "every path onto it.",
        "The \"% of a line in each class\" weighting falls out automatically: "
        "the router sums cost along the path, so a route 30% medium / 70% "
        "none pays 0.3×1.3 + 0.7×1.0 = 1.09× — never computed explicitly.",
        "Ground outside the risk region is written as nodata, so no route can "
        "leave the region and dodge the penalty entirely.",
    ], None),

    ("How the penalty acts on the cost surface",
     [], os.path.join(COST_FUT, "cost_penalty_explainer_zoom.png")),

    ("What the GeoPackage fields count",
     [], os.path.join(COST_FUT, "cost_penalty_explainer_cell.png")),

    ("5. Result — all 1,509 existing TVA routes re-costed", [
        "Every existing routed line re-priced on the penalised surface, on the "
        "routing pipeline's own cost basis:",
        "        system total   1.095×10¹² → 1.315×10¹²   =  +20.17%",
        "        median route   +17.68%      ·   upper quartile  +28.93%",
        "        1,493 of 1,509 routes increased  ·  16 unaffected  ·  none "
        "cheaper",
        "Projecting the fire weather forward roughly DOUBLES the penalty. On "
        "the historical classes the same method gives +10.74% system-wide and "
        "+7.98% for the median route, with 148 routes untouched — against only "
        "16 untouched here.",
        "This is \"same path, new price\": an upper bound, because these routes "
        "were optimised against the unpenalised surface. A re-run router would "
        "detour around penalised ground and pay less.",
    ], None),

    ("Cost before and after the penalty, per route",
     [], os.path.join(COST_FUT, "tva_route_cost_summary.png")),

    ("6. Caveats", [
        "FIRE WEATHER ONLY — FWI has no fuel term. It cannot see that the "
        "eastern TVA plateau carries far more fuel than the agricultural "
        "west, so the risk zones still concentrate in the west.",
        "    · fuel and fire weather are near-independent here (rank "
        "correlation −0.04), and fuel runs the OTHER way across TVA (+0.21 "
        "with longitude vs −0.52 for fire weather)",
        "    · LANDFIRE will be brought in to define the risk zones with fuel "
        "included. THESE NUMBERS ARE PROVISIONAL and will move — the workflow "
        "is what is being shared now, not the final values.",
        "Classes are region-relative: TVA \"high\" (FWI ≥ 40.8) is far milder "
        "than SoCal \"high\" (≥ 125.2). Not comparable across regions.",
        "One scenario, one realisation (ssp245, r1i1p1f1) — no ensemble "
        "spread. A 35-year future window is compared against a 15-year "
        "historical one.",
        "Not every cell worsens: 12.2% of TVA cells have a lower future p98 "
        "than historical. Worth checking whether that is real regional "
        "response or single-member noise.",
        "The penalty is blocky at 2.8 km — the fire-weather grid is far "
        "coarser than the 90 m routing grid.",
    ], None),
]


def add_picture_fitted(slide, path, top_in, max_h_in, slide_w_in=13.333):
    """Add a picture centred horizontally and scaled to fit the box."""
    pic = slide.shapes.add_picture(path, Inches(0), Inches(top_in),
                                   width=Inches(slide_w_in - 1.0))
    if pic.height > Inches(max_h_in):
        scale = Inches(max_h_in) / pic.height
        pic.height = int(pic.height * scale)
        pic.width = int(pic.width * scale)
    pic.left = int((Inches(slide_w_in) - pic.width) / 2)
    return pic


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)

    s = prs.slides.add_slide(prs.slide_layouts[0])
    s.shapes.title.text = TITLE
    s.placeholders[1].text = SUBTITLE

    for title, bullets, img in SLIDES:
        layout = prs.slide_layouts[5] if img else prs.slide_layouts[1]
        s = prs.slides.add_slide(layout)
        s.shapes.title.text = title
        s.shapes.title.text_frame.paragraphs[0].font.size = Pt(28)
        # The stock layouts carry 4:3 geometry (9 in wide); widening the slide
        # does not move them, so stretch the placeholders to the real width or
        # the text sits in the left two-thirds.
        for sh in s.shapes:
            if sh.is_placeholder:
                sh.left, sh.width = Inches(0.55), Inches(12.2)

        if img:
            add_picture_fitted(s, img, top_in=1.35, max_h_in=5.75)
        else:
            tf = s.placeholders[1].text_frame
            tf.word_wrap = True
            for i, b in enumerate(bullets):
                p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
                p.text = b
                # indented continuation lines are set smaller and unbulleted
                p.font.size = Pt(15) if b.startswith(" ") else Pt(17)
                p.level = 1 if b.startswith(" ") else 0

    os.makedirs(COST_FUT, exist_ok=True)
    prs.save(OUT)
    print(f"  saved -> {OUT}  ({len(prs.slides.__iter__.__self__._sldIdLst)} "
          f"slides)")


if __name__ == "__main__":
    main()
